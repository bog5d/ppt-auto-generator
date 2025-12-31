#!/usr/bin/env python3
"""
PPT自动生成器 v4.0 - 模板版
核心改进：
1. 先读JSON再下载图片（使用JSON中的提示词）
2. 多主题支持（4种预设+自定义）
3. 图片路径智能同步
4. 金句智能避让
5. 完整提示词显示
6. 自动生成AI提示词（基于image_desc）
7. 【新】支持企业模板导入（提取样式自动生成）
8. 【新】模板样式分析报告

作者：AI资源指挥官
版本：4.0
更新：2025-12-31
"""

import json
import sys
import os
import requests
import time
from datetime import datetime

# GUI支持
try:
    import tkinter as tk
    from tkinter import scrolledtext, messagebox, filedialog
    HAS_TKINTER = True
except ImportError:
    HAS_TKINTER = False

# 模板解析模块（可选）
try:
    from template_parser import (
        TemplateStyleExtractor, 
        TemplateBasedGenerator,
        analyze_template,
        get_theme_from_template,
        generate_from_template
    )
    HAS_TEMPLATE_PARSER = True
except ImportError:
    HAS_TEMPLATE_PARSER = False
    print("💡 提示：如需使用模板功能，请确保 template_parser.py 在同一目录下")

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE


# ========================================================================
# 核心类：PPT生成器 v3.8
# ========================================================================

class AutoPPTGeneratorV3:
    """智能PPT生成器 v3.8"""
    
    # 多主题配置
    THEMES = {
        'military_solemn': {
            'name': '军事庄重',
            'primary': RGBColor(26, 35, 126),      # 深蓝
            'accent': RGBColor(213, 0, 0),         # 中国红
            'text': RGBColor(33, 33, 33),          # 深灰文字
            'bg': RGBColor(250, 250, 250),         # 浅灰背景
            'quote': RGBColor(0, 150, 136),        # 青色（金句）
            'chart': RGBColor(63, 81, 181),        # 靛蓝（图表）
        },
        'tech_blue': {
            'name': '科技蓝',
            'primary': RGBColor(0, 119, 200),      # 科技蓝
            'accent': RGBColor(255, 152, 0),       # 橙色
            'text': RGBColor(33, 33, 33),
            'bg': RGBColor(250, 250, 250),
            'quote': RGBColor(0, 150, 136),
            'chart': RGBColor(0, 119, 200),
        },
        'nature_green': {
            'name': '自然绿',
            'primary': RGBColor(46, 125, 50),      # 自然绿
            'accent': RGBColor(255, 193, 7),       # 金黄
            'text': RGBColor(33, 33, 33),
            'bg': RGBColor(250, 250, 250),
            'quote': RGBColor(0, 121, 107),
            'chart': RGBColor(46, 125, 50),
        },
        'business_gray': {
            'name': '商务灰',
            'primary': RGBColor(66, 66, 66),       # 商务灰
            'accent': RGBColor(0, 150, 136),       # 青绿
            'text': RGBColor(33, 33, 33),
            'bg': RGBColor(250, 250, 250),
            'quote': RGBColor(0, 121, 107),
            'chart': RGBColor(96, 125, 139),
        }
    }
    
    # 6种智能布局
    LAYOUTS = {
        'left_text_right_image': {
            'name': '左文右图',
            'text_area': (0.3, 1.3, 4.5, 3.5),
            'image_area': (5.0, 1.3, 4.5, 3.5)
        },
        'right_text_left_image': {
            'name': '右文左图',
            'text_area': (5.0, 1.3, 4.5, 3.5),
            'image_area': (0.3, 1.3, 4.5, 3.5)
        },
        'top_text_bottom_image': {
            'name': '上文下图',
            'text_area': (0.3, 1.2, 9.4, 1.5),
            'image_area': (2.5, 2.8, 5, 2.2)  # 缩小图片，避免与金句重叠
        },
        'large_image_small_text': {
            'name': '大图配文',
            'image_area': (0.3, 1.2, 5.5, 3.5),
            'text_area': (6.0, 1.3, 3.5, 3.5)
        },
        'balanced': {
            'name': '平衡布局',
            'text_area': (0.3, 1.3, 4.5, 3.5),
            'image_area': (5.0, 1.3, 4.5, 3.5)
        },
        'emphasis_text': {
            'name': '文字为主',
            'text_area': (0.3, 1.3, 6.2, 3.5),
            'image_area': (6.8, 1.5, 2.8, 3)
        }
    }
    
    def __init__(self, theme='military_solemn'):
        """初始化生成器"""
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(5.625)
        self.theme = self.THEMES.get(theme, self.THEMES['military_solemn'])
        self.slide_index = 0
        
        print(f"🎨 使用主题: {self.theme.get('name', theme)}")
    
    def generate_from_json(self, json_path_or_data, output_path):
        """从JSON生成完整PPT (支持文件路径或直接传入数据)"""
        if isinstance(json_path_or_data, dict):
            # 直接传入的JSON数据
            data = json_path_or_data
        else:
            # 从文件读取
            with open(json_path_or_data, 'r', encoding='utf-8') as f:
                data = json.load(f)
        
        metadata = data.get('metadata', {})
        slides_data = data.get('slides', [])
        
        print(f"\n{'='*60}")
        print(f"🚀 开始生成 PPT...")
        print(f"{'='*60}\n")
        
        for slide_data in slides_data:
            slide_type = slide_data.get('type')
            
            if slide_type == 'cover':
                self.create_cover_slide(slide_data)
            elif slide_type == 'section':
                self.create_section_slide(slide_data)
            elif slide_type == 'content_image':
                self.create_content_with_image_slide(slide_data)
            elif slide_type == 'chart':
                self.create_chart_slide(slide_data)
            elif slide_type == 'ending':
                self.create_ending_slide(slide_data)
        
        self.prs.save(output_path)
        
        print(f"\n{'='*60}")
        print(f"✅ PPT生成成功！")
        print(f"{'='*60}")
        print(f"📊 总页数: {len(self.prs.slides)}")
        print(f"🎨 主题: {self.theme.get('name', 'default')}")
        print(f"📁 输出路径: {output_path}")
        print(f"{'='*60}\n")
    
    def auto_select_layout(self, data):
        """智能选择布局（循环切换）"""
        layouts = list(self.LAYOUTS.keys())
        return layouts[self.slide_index % len(layouts)]
    
    def add_structured_bullets(self, text_frame, bullets):
        """添加结构化文字（支持"标题：内容"格式）- 智能换行和字号"""
        text_frame.word_wrap = True
        
        # 更小的字体大小，避免溢出
        title_size = 9
        content_size = 8
        
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = text_frame.paragraphs[0]
                p.text = ''  # 清空默认文本
            else:
                p = text_frame.add_paragraph()
            
            # 检测"标题：内容"格式
            if '：' in bullet or ':' in bullet:
                parts = bullet.split('：', 1) if '：' in bullet else bullet.split(':', 1)
                
                if len(parts) == 2:
                    title_text = parts[0].strip()
                    content_text = parts[1].strip()
                    
                    # 标题部分（加粗）
                    run1 = p.add_run()
                    run1.text = title_text + '：'
                    run1.font.bold = True
                    run1.font.size = Pt(title_size)
                    run1.font.color.rgb = self.theme['primary']
                    
                    # 如果内容过长（超过25字），强制换行到新段落
                    if len(content_text) > 25:
                        # 添加换行，内容放在下一行
                        p2 = text_frame.add_paragraph()
                        run2 = p2.add_run()
                        run2.text = '  ' + content_text  # 缩进
                        run2.font.size = Pt(content_size)
                        run2.font.color.rgb = self.theme['text']
                        p2.level = 0
                        p2.space_before = Pt(0)
                        p2.space_after = Pt(2)
                        p2.line_spacing = 1.05
                    else:
                        # 内容部分（普通，同一行）
                        run2 = p.add_run()
                        run2.text = content_text
                        run2.font.size = Pt(content_size)
                        run2.font.color.rgb = self.theme['text']
                else:
                    # 普通文字
                    run = p.add_run()
                    run.text = bullet
                    run.font.size = Pt(title_size)
                    run.font.color.rgb = self.theme['text']
            else:
                # 普通文字 - 超长也要换行
                if len(bullet) > 35:
                    # 分割成多行
                    words = bullet
                    while len(words) > 35:
                        # 找到合适的分割点
                        split_pos = 35
                        for punct in ['，', '、', '；', '。', ' ']:
                            pos = words[:40].rfind(punct)
                            if pos > 20:
                                split_pos = pos + 1
                                break
                        
                        run = p.add_run()
                        run.text = words[:split_pos]
                        run.font.size = Pt(content_size)
                        run.font.color.rgb = self.theme['text']
                        
                        words = words[split_pos:]
                        if words:
                            p = text_frame.add_paragraph()
                            p.level = 0
                            p.space_before = Pt(0)
                            p.space_after = Pt(1)
                    
                    if words:
                        run = p.add_run()
                        run.text = words
                        run.font.size = Pt(content_size)
                        run.font.color.rgb = self.theme['text']
                else:
                    run = p.add_run()
                    run.text = bullet
                    run.font.size = Pt(title_size)
                    run.font.color.rgb = self.theme['text']
            
            p.level = 0
            p.space_before = Pt(1)
            p.space_after = Pt(1)
            p.line_spacing = 1.05
    
    def create_cover_slide(self, data):
        """封面页"""
        layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(layout)
        
        # 背景
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.theme['primary']
        
        # 主标题 - 自动调整字号
        title_text = data.get('title', '')
        title_box = slide.shapes.add_textbox(
            Inches(0.3), Inches(1.5), Inches(9.4), Inches(1.5)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.text = title_text
        
        p = tf.paragraphs[0]
        # 根据标题长度自动调整字号
        if len(title_text) > 20:
            p.font.size = Pt(32)
        elif len(title_text) > 15:
            p.font.size = Pt(36)
        else:
            p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # 副标题
        if data.get('subtitle'):
            subtitle_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(3.2), Inches(9), Inches(0.8)
            )
            tf = subtitle_box.text_frame
            tf.word_wrap = True
            tf.text = data['subtitle']
            
            p = tf.paragraphs[0]
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(230, 230, 230)
            p.alignment = PP_ALIGN.CENTER
        
        # 口号
        if data.get('slogan'):
            slogan_box = slide.shapes.add_textbox(
                Inches(2), Inches(4.5), Inches(6), Inches(0.6)
            )
            tf = slogan_box.text_frame
            tf.text = data['slogan']
            
            p = tf.paragraphs[0]
            p.font.size = Pt(14)
            p.font.italic = True
            p.font.color.rgb = self.theme['accent']
            p.alignment = PP_ALIGN.CENTER
        
        self.slide_index += 1
        return slide
    
    def create_section_slide(self, data):
        """章节页"""
        layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(layout)
        
        # 背景
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.theme['bg']
        
        # 装饰条
        deco = slide.shapes.add_shape(
            1, Inches(0), Inches(2.3),
            Inches(10), Inches(1)
        )
        deco.fill.solid()
        deco.fill.fore_color.rgb = self.theme['primary']
        deco.line.color.rgb = self.theme['primary']
        
        # 标题 - 自动调整字号
        title_text = data.get('title', '')
        title_box = slide.shapes.add_textbox(
            Inches(0.3), Inches(2.3), Inches(9.4), Inches(1)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.text = title_text
        
        p = tf.paragraphs[0]
        # 根据标题长度自动调整字号
        if len(title_text) > 16:
            p.font.size = Pt(32)
        elif len(title_text) > 12:
            p.font.size = Pt(38)
        else:
            p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        self.slide_index += 1
        return slide
    
    def create_content_with_image_slide(self, data):
        """图文混排页（智能布局）"""
        layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(layout)
        
        # 背景
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.theme['bg']
        
        # 标题 - 自动调整字号
        title_text = data.get('title', '')
        title_box = slide.shapes.add_textbox(
            Inches(0.3), Inches(0.3), Inches(9.4), Inches(0.8)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.text = title_text
        
        p = tf.paragraphs[0]
        # 根据标题长度自动调整字号
        if len(title_text) > 18:
            p.font.size = Pt(24)
        elif len(title_text) > 12:
            p.font.size = Pt(28)
        else:
            p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.theme['primary']
        
        # 标题下划线
        line = slide.shapes.add_shape(
            1, Inches(0.3), Inches(1.1),
            Inches(2), Inches(0)
        )
        line.line.color.rgb = self.theme['accent']
        line.line.width = Pt(3)
        
        # 智能选择布局
        layout_type = data.get('layout', self.auto_select_layout(data))
        layout_config = self.LAYOUTS[layout_type]
        
        print(f"  → 第{self.slide_index+1}页使用布局: {layout_config['name']}")
        
        # 文字区域
        text_area = layout_config['text_area']
        content_box = slide.shapes.add_textbox(
            Inches(text_area[0]), Inches(text_area[1]),
            Inches(text_area[2]), Inches(text_area[3])
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        
        bullets = data.get('bullets', [])
        self.add_structured_bullets(tf, bullets)
        
        # 图片区域
        image_area = layout_config['image_area']
        image_path = data.get('image')
        image_prompt = data.get('image_prompt', '')
        
        # 检查图片路径是否存在（尝试多个可能的路径）
        image_exists = False
        actual_path = None
        
        if image_path:
            # 尝试的路径列表
            possible_paths = [
                image_path,                                    # 原始路径
                os.path.basename(image_path),                  # 当前目录
                os.path.join(os.getcwd(), os.path.basename(image_path)),  # 工作目录
                os.path.join('C:\\Users\\王波', os.path.basename(image_path)),  # 用户目录
            ]
            
            for p in possible_paths:
                if os.path.exists(p):
                    actual_path = p
                    image_exists = True
                    break
        
        if image_exists and actual_path:
            try:
                print(f"  📷 插入图片: {os.path.basename(actual_path)}")
                slide.shapes.add_picture(
                    actual_path,  # 使用找到的实际路径
                    Inches(image_area[0]), Inches(image_area[1]),
                    width=Inches(image_area[2]),
                    height=Inches(image_area[3])
                )
                # 不再显示提示词（避免与金句重叠）
            except Exception as e:
                print(f"  ⚠️ 图片插入失败: {e}")
                self._add_image_placeholder(
                    slide, 
                    data.get('image_desc', '图片'),
                    image_area,
                    None  # 不传提示词
                )
        else:
            if image_path:
                print(f"  ⚠️ 图片不存在: {image_path}")
            self._add_image_placeholder(
                slide,
                data.get('image_desc', '图片'),
                image_area,
                None  # 不传提示词，避免与金句重叠
            )
        
        # 金句（智能避让） - 放在页面底部固定位置，不与图片重叠
        if data.get('quote') and layout_type not in ['large_image_small_text']:
            # 金句固定在页面最底部
            quote_y = 5.15
            
            quote_box = slide.shapes.add_textbox(
                Inches(0.3), Inches(quote_y), 
                Inches(9.4), Inches(0.4)
            )
            tf = quote_box.text_frame
            tf.word_wrap = True
            
            # 截断过长的金句
            quote_text = data["quote"]
            if len(quote_text) > 60:
                quote_text = quote_text[:57] + '...'
            tf.text = f'💡 {quote_text}'
            tf.word_wrap = True
            
            p = tf.paragraphs[0]
            p.font.size = Pt(12)
            p.font.italic = True
            p.font.color.rgb = self.theme['quote']
        
        self.slide_index += 1
        return slide
    
    def _add_prompt_text(self, slide, area_config, prompt):
        """在图片下方显示生成提示词"""
        prompt_box = slide.shapes.add_textbox(
            Inches(area_config[0]),
            Inches(area_config[1] + area_config[3] + 0.05),
            Inches(area_config[2]),
            Inches(0.4)
        )
        tf = prompt_box.text_frame
        tf.text = f"Prompt: {prompt}"
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.font.size = Pt(8)
        p.font.italic = True
        p.font.color.rgb = RGBColor(100, 100, 100)
    
    def _add_image_placeholder(self, slide, description, area_config, prompt=None):
        """添加专业图片占位符 + 显示生成提示词"""
        # 背景框
        placeholder = slide.shapes.add_shape(
            1,
            Inches(area_config[0]), Inches(area_config[1]),
            Inches(area_config[2]), Inches(area_config[3])
        )
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor(245, 248, 250)
        placeholder.line.color.rgb = self.theme['primary']
        placeholder.line.width = Pt(2)
        placeholder.line.dash_style = 2
        
        # 图标框
        icon_size = 0.5
        icon_box = slide.shapes.add_shape(
            1,
            Inches(area_config[0] + area_config[2]/2 - icon_size/2),
            Inches(area_config[1] + area_config[3]/2 - icon_size - 0.3),
            Inches(icon_size), Inches(icon_size)
        )
        icon_box.fill.solid()
        icon_box.fill.fore_color.rgb = self.theme['primary']
        icon_box.line.color.rgb = self.theme['primary']
        
        # 图标文字
        icon_text = slide.shapes.add_textbox(
            Inches(area_config[0] + area_config[2]/2 - icon_size/2),
            Inches(area_config[1] + area_config[3]/2 - icon_size - 0.3),
            Inches(icon_size), Inches(icon_size)
        )
        tf = icon_text.text_frame
        tf.text = "🖼️"
        p = tf.paragraphs[0]
        p.font.size = Pt(28)
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # 描述文字
        text_box = slide.shapes.add_textbox(
            Inches(area_config[0] + 0.3), 
            Inches(area_config[1] + area_config[3]/2 + 0.1),
            Inches(area_config[2] - 0.6), Inches(0.8)
        )
        tf = text_box.text_frame
        tf.text = description
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = self.theme['primary']
        p.alignment = PP_ALIGN.CENTER
        
        # 提示文字 - 显示实际提示词（如果有）
        hint_text = prompt if prompt else "(参考提示词替换图片)"
        hint_box = slide.shapes.add_textbox(
            Inches(area_config[0] + 0.3),
            Inches(area_config[1] + area_config[3]/2 + 0.7),
            Inches(area_config[2] - 0.6), Inches(0.6)
        )
        tf = hint_box.text_frame
        tf.text = hint_text
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.font.size = Pt(8)
        p.font.italic = True
        p.font.color.rgb = RGBColor(80, 120, 160) if prompt else RGBColor(120, 120, 120)
        p.alignment = PP_ALIGN.CENTER
    
    def create_chart_slide(self, data):
        """图表页"""
        layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(layout)
        
        # 背景
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.theme['bg']
        
        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4), Inches(9), Inches(0.6)
        )
        tf = title_box.text_frame
        tf.text = data.get('title', '')
        
        p = tf.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.theme['primary']
        
        # 图表
        chart_data_config = data.get('chart_data', {})
        chart_type = data.get('chart_type', 'column')
        
        chart_data = CategoryChartData()
        chart_data.categories = chart_data_config.get('labels', [])
        
        for dataset in chart_data_config.get('datasets', []):
            chart_data.add_series(dataset['name'], dataset['values'])
        
        x, y, cx, cy = Inches(1.5), Inches(1.5), Inches(7), Inches(3.5)
        
        if chart_type == 'column':
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
            ).chart
        
        # 备注
        if data.get('note'):
            note_box = slide.shapes.add_textbox(
                Inches(1), Inches(5.1), Inches(8), Inches(0.4)
            )
            tf = note_box.text_frame
            tf.text = data['note']
            
            p = tf.paragraphs[0]
            p.font.size = Pt(10)
            p.font.italic = True
            p.font.color.rgb = RGBColor(120, 120, 120)
        
        self.slide_index += 1
        return slide
    
    def create_ending_slide(self, data):
        """结束页"""
        layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(layout)
        
        # 背景
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.theme['bg']
        
        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.6), Inches(9), Inches(0.8)
        )
        tf = title_box.text_frame
        tf.text = data.get('title', '')
        
        p = tf.paragraphs[0]
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = self.theme['primary']
        
        # 要点列表
        if data.get('bullets'):
            content_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(1.6), Inches(7), Inches(2.8)
            )
            tf = content_box.text_frame
            tf.word_wrap = True
            
            self.add_structured_bullets(tf, data['bullets'])
        
        # 金句
        if data.get('quote'):
            quote_box = slide.shapes.add_textbox(
                Inches(1), Inches(4.6), Inches(8), Inches(0.8)
            )
            tf = quote_box.text_frame
            tf.text = f'💡 {data["quote"]}'
            
            p = tf.paragraphs[0]
            p.font.size = Pt(16)
            p.font.italic = True
            p.font.bold = True
            p.font.color.rgb = self.theme['accent']
            p.alignment = PP_ALIGN.CENTER
        
        self.slide_index += 1
        return slide


# ========================================================================
# 图片下载模块
# ========================================================================

def download_single_image_unsplash(query, api_key, filename):
    """从Unsplash下载单张图片"""
    try:
        url = "https://api.unsplash.com/search/photos"
        params = {
            'query': query,
            'client_id': api_key,
            'per_page': 1,
            'orientation': 'landscape'
        }
        
        response = requests.get(url, params=params, timeout=15)
        
        if response.status_code == 200:
            data = response.json()
            if data.get('results'):
                img_url = data['results'][0]['urls']['regular']
                photographer = data['results'][0]['user']['name']
                
                img_response = requests.get(img_url, timeout=15)
                if img_response.status_code == 200:
                    with open(filename, 'wb') as f:
                        f.write(img_response.content)
                    return True, f"成功 (摄影师: {photographer})"
        
        return False, f"API返回错误: {response.status_code}"
    
    except requests.exceptions.Timeout:
        return False, "请求超时"
    except Exception as e:
        return False, f"异常: {str(e)}"


def generate_single_image_siliconflow(prompt, api_key, filename, max_retries=3):
    """使用硅基流动API生成单张图片（带重试机制）"""
    
    for attempt in range(max_retries):
        try:
            url = "https://api.siliconflow.cn/v1/images/generations"
            headers = {
                "Authorization": f"Bearer {api_key}",
                "Content-Type": "application/json"
            }
            data = {
                "model": "black-forest-labs/FLUX.1-schnell",
                "prompt": prompt,
                "image_size": "1024x1024",
                "num_inference_steps": 20
            }
            
            response = requests.post(url, headers=headers, json=data, timeout=120)
            
            if response.status_code == 200:
                result = response.json()
                
                if 'images' in result and len(result['images']) > 0:
                    img_data = result['images'][0]
                    if 'url' in img_data:
                        img_response = requests.get(img_data['url'], timeout=15)
                        with open(filename, 'wb') as f:
                            f.write(img_response.content)
                        return True, "AI生成成功"
                    elif 'b64_json' in img_data:
                        import base64
                        img_bytes = base64.b64decode(img_data['b64_json'])
                        with open(filename, 'wb') as f:
                            f.write(img_bytes)
                        return True, "AI生成成功"
                
                return False, "返回格式不支持"
            
            elif response.status_code == 429:
                # API限流，等待后重试
                wait_time = 5 * (attempt + 1)  # 5秒、10秒、15秒
                if attempt < max_retries - 1:
                    print(f"  ⏳ API限流，等待{wait_time}秒后重试...")
                    time.sleep(wait_time)
                    continue
                return False, "API限流，重试失败"
            
            return False, f"API返回错误: {response.status_code}"
        
        except requests.exceptions.Timeout:
            if attempt < max_retries - 1:
                print(f"  ⏳ 请求超时，重试中...")
                time.sleep(3)
                continue
            return False, "生成超时"
        except Exception as e:
            return False, f"异常: {str(e)}"
    
    return False, "重试次数用尽"


def generate_smart_prompt(title, bullets, desc):
    """
    根据幻灯片标题和内容智能生成高质量AI图片提示词
    使用英文提示词以获得更好的AI生成效果
    """
    # 提取bullets中的关键词
    keywords = []
    if bullets:
        for bullet in bullets[:3]:  # 只取前3个要点
            # 提取冒号前的关键词（通常是标题）
            if '：' in bullet:
                key = bullet.split('：')[0].strip()
                keywords.append(key)
            elif ':' in bullet:
                key = bullet.split(':')[0].strip()
                keywords.append(key)
            else:
                # 取前10个字作为关键词
                keywords.append(bullet[:10])
    
    # 关键词映射表：中文主题 -> 英文描述
    keyword_mapping = {
        # 电磁/军事相关
        '电磁': 'electromagnetic waves, radar systems',
        '雷达': 'military radar system, antenna array',
        '脉冲': 'electromagnetic pulse, EMP effect',
        '攻击': 'cyber attack visualization, security threat',
        '防护': 'protective shield, defense system',
        '辐射': 'radiation protection, electromagnetic shielding',
        '屏蔽': 'metal shielding box, Faraday cage',
        '干扰': 'electronic jamming, signal interference',
        '通信': 'communication systems, satellite links',
        '导弹': 'missile defense system, military technology',
        '战场': 'modern battlefield, military operations',
        '武器': 'advanced weapons system, military equipment',
        '传导': 'electrical conduction, circuit protection',
        '耦合': 'electromagnetic coupling, signal transmission',
        '滤波': 'electronic filter, signal processing',
        '芯片': 'microchip, semiconductor technology',
        '设备': 'electronic equipment, technical devices',
        '系统': 'integrated system, technical architecture',
        '标准': 'technical standards, certification documents',
        '试验': 'laboratory testing, scientific experiment',
        '验证': 'verification process, quality control',
        # 通用技术
        '技术': 'advanced technology, innovation',
        '科技': 'high-tech, futuristic design',
        '数据': 'data visualization, digital information',
        '网络': 'network topology, cyber infrastructure',
        '安全': 'security systems, protection measures',
    }
    
    # 构建英文描述
    english_elements = []
    
    # 从标题和关键词中匹配
    all_text = title + ' ' + ' '.join(keywords)
    for cn_key, en_value in keyword_mapping.items():
        if cn_key in all_text:
            english_elements.append(en_value)
    
    # 如果没有匹配到，使用默认描述
    if not english_elements:
        english_elements = ['technical illustration', 'professional diagram']
    
    # 限制元素数量
    english_elements = english_elements[:3]
    
    # 构建完整提示词
    subject = ', '.join(english_elements)
    
    # 专业的AI图片生成提示词模板
    prompt = f"""Professional technical illustration showing {subject}. 
Style: Clean modern infographic, technical diagram, blueprint aesthetic.
Colors: Blue and white color scheme, professional look.
Quality: High resolution, 4K, detailed, sharp focus.
Background: Clean gradient or solid color, minimalist.
NO text, NO watermarks, NO human faces."""
    
    # 压缩为一行
    prompt = ' '.join(prompt.split())
    
    return prompt


def get_text_from_gui():
    """
    弹出GUI窗口让用户输入/粘贴大纲文本
    返回用户输入的文本，如果取消则返回None
    """
    if not HAS_TKINTER:
        print("⚠️  未安装tkinter，使用终端输入模式")
        return None
    
    result = {'text': None}
    
    def on_submit():
        result['text'] = text_area.get('1.0', tk.END).strip()
        if result['text']:
            root.destroy()
        else:
            messagebox.showwarning("提示", "请输入大纲内容")
    
    def on_cancel():
        result['text'] = None
        root.destroy()
    
    # 创建窗口 - 更大尺寸
    root = tk.Tk()
    root.title("📝 PPT大纲输入 - 粘贴您的大纲文本")
    root.geometry("1000x750")
    root.configure(bg='#2b2b2b')
    
    # 说明标签 - 更大字体
    label = tk.Label(root, 
                     text="📋 请粘贴大纲文本（支持多种格式，会自动智能解析）",
                     font=('Microsoft YaHei', 14, 'bold'), bg='#2b2b2b', fg='#ffffff')
    label.pack(pady=15)
    
    # 格式提示
    hint = tk.Label(root, 
                    text="支持格式：# 标题 | ## 章节/第X页 | ### 内容 | - 要点 | > 金句",
                    font=('Microsoft YaHei', 11), bg='#2b2b2b', fg='#aaaaaa')
    hint.pack(pady=5)
    
    # 按钮框架 - 先pack按钮，固定在底部
    btn_frame = tk.Frame(root, bg='#2b2b2b')
    btn_frame.pack(side=tk.BOTTOM, pady=20, fill=tk.X)
    
    submit_btn = tk.Button(btn_frame, text="✅ 确认生成PPT", command=on_submit,
                          font=('Microsoft YaHei', 16, 'bold'), bg='#4CAF50', fg='white',
                          width=20, height=2, cursor='hand2', relief='raised', bd=3)
    submit_btn.pack(side=tk.LEFT, padx=50, expand=True)
    
    cancel_btn = tk.Button(btn_frame, text="❌ 取消", command=on_cancel,
                          font=('Microsoft YaHei', 16), bg='#f44336', fg='white',
                          width=14, height=2, cursor='hand2', relief='raised', bd=3)
    cancel_btn.pack(side=tk.LEFT, padx=50, expand=True)
    
    # 快捷键提示
    shortcut_label = tk.Label(btn_frame, text="💡 快捷键: Ctrl+Enter 确认",
                              font=('Microsoft YaHei', 10), bg='#2b2b2b', fg='#888888')
    shortcut_label.pack(side=tk.RIGHT, padx=20)
    
    # 文本输入区 - 更大字体
    text_area = scrolledtext.ScrolledText(root, width=100, height=25, 
                                          font=('Consolas', 13),
                                          wrap=tk.WORD, bg='#1e1e1e', fg='#d4d4d4',
                                          insertbackground='white')
    text_area.pack(padx=30, pady=10, fill=tk.BOTH, expand=True)
    
    # 绑定 Ctrl+Enter 快捷键
    def on_ctrl_enter(event):
        on_submit()
        return 'break'
    text_area.bind('<Control-Return>', on_ctrl_enter)
    
    # 居中显示
    root.update_idletasks()
    x = (root.winfo_screenwidth() - root.winfo_width()) // 2
    y = (root.winfo_screenheight() - root.winfo_height()) // 2
    root.geometry(f"+{x}+{y}")
    
    # 让文本框获得焦点
    text_area.focus_set()
    
    root.mainloop()
    
    return result['text']


def parse_outline_to_json(text):
    """
    智能解析大纲文本转换为JSON结构
    支持多种格式：
    # 标题 / # XXX大纲         -> 提取标题
    ## 第X页：标题 / ## 章节名  -> section 或 content_image
    ### 内容标题              -> content_image
    - **标题**：内容 / - 内容  -> bullets
    > 金句                    -> quote
    ---                       -> 分隔符（忽略）
    """
    import re
    
    lines = text.strip().split('\n')
    
    slides = []
    current_slide = None
    cover_title = "演示文稿"
    cover_subtitle = ""
    layout_index = 0
    layouts = ['left_text_right_image', 'right_text_left_image']
    
    def clean_text(s):
        """清理Markdown格式"""
        s = s.replace('**', '').replace('*', '')
        s = re.sub(r'\[.*?\]', '', s)  # 移除[内容]
        s = s.replace('"', '').replace('"', '').replace('"', '')
        return s.strip()
    
    def extract_title_from_section(line):
        """从 '## 第X页：标题' 或 '## 标题' 提取标题"""
        line = line.lstrip('#').strip()
        # 匹配 "第X页：标题" 或 "第X页:标题"
        match = re.match(r'第.+[页节][\s：:]+(.+)', line)
        if match:
            return clean_text(match.group(1))
        # 匹配 "封面" "总结" 等特殊页
        if '封面' in line or '结语' in line or '结尾' in line:
            return None  # 封面/结尾特殊处理
        return clean_text(line)
    
    def is_section_only(title):
        """判断是否只是章节页（没有具体内容的）"""
        keywords = ['目录', '章节', '第一部分', '第二部分', '第三部分', '第四部分', '概述', '引言']
        return any(kw in title for kw in keywords)
    
    for line in lines:
        original_line = line
        line = line.strip()
        
        # 跳过空行和分隔符
        if not line or line == '---' or line.startswith('---'):
            continue
        
        # 主标题 (# 开头，但不是##)
        if line.startswith('# ') and not line.startswith('## '):
            title = line[2:].strip()
            # 处理 "# XXX大纲" 格式
            title = re.sub(r'大纲$', '', title).strip()
            title = clean_text(title)
            if title:
                cover_title = title
            continue
        
        # 二级标题 (## 开头)
        if line.startswith('## ') and not line.startswith('### '):
            # 保存上一个slide
            if current_slide:
                slides.append(current_slide)
            
            section_title = extract_title_from_section(line)
            
            # 检查是否是封面
            if section_title is None or '封面' in line:
                current_slide = {'type': '_cover_placeholder'}
                continue
            
            # 检查是否是总结/结尾
            if '总结' in section_title or '结语' in section_title or '结尾' in section_title:
                current_slide = {
                    'type': 'ending',
                    'title': section_title,
                    'bullets': [],
                    'quote': ''
                }
                continue
            
            # 普通内容页（带图片）
            current_slide = {
                'type': 'content_image',
                'title': section_title,
                'bullets': [],
                'layout': layouts[layout_index % 2],
                'image_desc': f'{section_title}示意图',
                'image': f'images/slide_{layout_index + 1}.png'
            }
            layout_index += 1
            continue
        
        # 三级标题 (### 开头)
        if line.startswith('### '):
            # 保存上一个slide
            if current_slide and current_slide.get('type') != '_cover_placeholder':
                slides.append(current_slide)
            
            title = clean_text(line[4:])
            current_slide = {
                'type': 'content_image',
                'title': title,
                'bullets': [],
                'layout': layouts[layout_index % 2],
                'image_desc': f'{title}示意图',
                'image': f'images/slide_{layout_index + 1}.png'
            }
            layout_index += 1
            continue
        
        # 要点 (- 或 * 开头，支持缩进)
        if re.match(r'^[\s]*[-*]\s+', original_line):
            bullet_text = re.sub(r'^[\s]*[-*]\s+', '', original_line)
            bullet_text = clean_text(bullet_text)
            
            if current_slide:
                if current_slide.get('type') == '_cover_placeholder':
                    # 封面的要点提取为副标题/演讲人等
                    if '标题' in bullet_text and '：' in bullet_text:
                        cover_title = bullet_text.split('：', 1)[1].strip()
                    elif '副标题' in bullet_text and '：' in bullet_text:
                        cover_subtitle = bullet_text.split('：', 1)[1].strip()
                elif 'bullets' in current_slide and bullet_text:
                    current_slide['bullets'].append(bullet_text)
            continue
        
        # 金句 (> 开头)
        if line.startswith('> '):
            if current_slide and current_slide.get('type') not in [None, '_cover_placeholder']:
                current_slide['quote'] = clean_text(line[2:])
            continue
        
        # 理解类比等特殊段落
        if line.startswith('**') and '**' in line[2:]:
            # 可能是 **理解类比**：内容
            continue
        
        # 其他文本
        if current_slide is None:
            if not cover_subtitle:
                cover_subtitle = clean_text(line)
        elif 'bullets' in current_slide:
            cleaned = clean_text(line)
            if cleaned and not cleaned.startswith('**'):
                current_slide['bullets'].append(cleaned)
    
    # 保存最后一个slide
    if current_slide and current_slide.get('type') not in [None, '_cover_placeholder']:
        slides.append(current_slide)
    
    # 过滤掉占位符
    slides = [s for s in slides if s.get('type') != '_cover_placeholder']
    
    # 添加封面
    cover_slide = {
        'type': 'cover',
        'title': cover_title,
        'subtitle': cover_subtitle or '专业培训课程',
        'slogan': ''
    }
    
    # 检查是否已有ending
    has_ending = any(s.get('type') == 'ending' for s in slides)
    
    # 如果没有ending，添加一个
    if not has_ending:
        ending_slide = {
            'type': 'ending',
            'title': '谢谢观看',
            'bullets': ['欢迎交流讨论'],
            'quote': '合规运作，价值创造'
        }
        slides.append(ending_slide)
    
    # 组装完整JSON
    result = {
        'metadata': {
            'title': cover_title,
            'theme': 'business_gray',
            'version': '3.9',
            'total_slides': len(slides) + 1
        },
        'slides': [cover_slide] + slides
    }
    
    return result


def extract_image_prompts_from_json(json_data):
    """从JSON中提取所有图片提示词和路径"""
    image_tasks = []
    
    for slide in json_data.get('slides', []):
        if slide.get('type') == 'content_image':
            # 优先使用image_prompt，如果没有则根据slide内容智能生成
            prompt = slide.get('image_prompt', '')
            desc = slide.get('image_desc', '图片')
            title = slide.get('title', '')
            bullets = slide.get('bullets', [])
            
            # 如果没有prompt，使用智能生成
            if not prompt:
                prompt = generate_smart_prompt(title, bullets, desc)
                print(f"  💡 智能提示词 [{title[:15]}...]: {prompt[:70]}...")
            
            # 只要有描述就添加任务（无论是否有原始prompt）
            if desc or prompt:
                # 生成默认文件名或使用指定路径（使用绝对路径）
                image_path = slide.get('image', f"image_{len(image_tasks)+1}.jpg")
                
                # 如果是相对路径，转换为绝对路径（在当前目录）
                if not os.path.isabs(image_path):
                    image_path = os.path.abspath(image_path)
                
                image_tasks.append({
                    'prompt': prompt,
                    'path': image_path,
                    'desc': desc,
                    'title': title
                })
    
    return image_tasks


def download_images_from_json(image_tasks, unsplash_key=None, siliconflow_key=None):
    """根据JSON中的任务列表下载图片"""
    if not image_tasks:
        print("ℹ️  JSON中没有图片提示词，跳过下载\n")
        return True
    
    print("\n" + "=" * 70)
    print(f"🚀 智能图片下载系统 - 基于JSON配置")
    print("=" * 70)
    print(f"📅 开始时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
    print(f"📊 待下载图片数: {len(image_tasks)}")
    print()
    
    stats = {
        'unsplash_success': 0,
        'ai_success': 0,
        'failed': 0,
        'details': []
    }
    
    for i, task in enumerate(image_tasks, 1):
        prompt = task['prompt']
        filepath = task['path']
        desc = task['desc']
        
        print(f"[{i}/{len(image_tasks)}] {os.path.basename(filepath)}")
        print(f"  📝 描述: {desc}")
        
        # 确保目录存在
        dir_path = os.path.dirname(filepath)
        if dir_path:
            os.makedirs(dir_path, exist_ok=True)
        
        # 优先使用AI生成（使用JSON中的prompt）
        if siliconflow_key and prompt:
            print("  🤖 使用AI生成图片...")
            print(f"  📝 提示词: {prompt[:60]}{'...' if len(prompt) > 60 else ''}")
            success, msg = generate_single_image_siliconflow(
                prompt,
                siliconflow_key,
                filepath
            )
            
            if success:
                print(f"  ✅ AI生成成功")
                stats['ai_success'] += 1
                stats['details'].append({
                    'file': filepath,
                    'source': 'SiliconFlow AI',
                    'status': 'success',
                    'prompt': prompt
                })
                time.sleep(2)
                continue
            else:
                print(f"  ⚠️  AI生成失败: {msg}")
        
        # AI失败，尝试Unsplash备用
        if unsplash_key:
            print("  🔄 尝试Unsplash备用...")
            success, msg = download_single_image_unsplash(
                desc,  # 使用描述作为搜索词
                unsplash_key,
                filepath
            )
            
            if success:
                print(f"  ✅ Unsplash {msg}")
                stats['unsplash_success'] += 1
                stats['details'].append({
                    'file': filepath,
                    'source': 'Unsplash',
                    'status': 'success'
                })
                time.sleep(1)
                continue
            else:
                print(f"  ⚠️  Unsplash失败: {msg}")
        
        # 都失败
        print(f"  ❌ 所有下载源都失败，将使用占位图")
        stats['failed'] += 1
        stats['details'].append({
            'file': filepath,
            'source': 'None',
            'status': 'failed'
        })
    
    # 生成报告
    print("\n" + "=" * 70)
    print("📊 下载报告")
    print("=" * 70)
    print(f"✅ Unsplash成功: {stats['unsplash_success']}")
    print(f"✅ AI生成成功: {stats['ai_success']}")
    print(f"❌ 失败（使用占位图）: {stats['failed']}")
    print()
    
    if stats['details']:
        print("详细结果:")
        for item in stats['details']:
            status_icon = "✅" if item['status'] == 'success' else "❌"
            print(f"  {status_icon} {os.path.basename(item['file'])} - {item['source']}")
    
    print("=" * 70)
    print()
    
    # 返回成功的路径列表
    success_paths = [item['file'] for item in stats['details'] if item['status'] == 'success']
    return success_paths


# ========================================================================
# 主函数 v3.9
# ========================================================================

def select_template_file():
    """弹出文件选择器选择PPT模板"""
    if not HAS_TKINTER:
        # 终端输入模式
        return input("请输入模板文件路径: ").strip()
    
    root = tk.Tk()
    root.withdraw()  # 隐藏主窗口
    
    file_path = filedialog.askopenfilename(
        title="选择PPT模板文件",
        filetypes=[
            ("PowerPoint文件", "*.pptx"),
            ("所有文件", "*.*")
        ]
    )
    
    root.destroy()
    return file_path


def main():
    """主函数 v4.0 - 支持模板导入"""
    print("=" * 70)
    print("PPT自动生成器 v4.0 - 模板版")
    print("=" * 70)
    print()
    print("📌 v4.0 新特性：")
    print("  ✅ 粘贴大纲文本自动支持AI图片生成")
    print("  ✅ GUI弹窗输入（更方便粘贴）")
    print("  ✅ 4种主题配色支持")
    print("  ✅ 图片路径智能同步")
    print("  🆕 支持企业模板导入（提取样式自动生成）")
    print("  🆕 模板分析报告功能")
    
    if HAS_TEMPLATE_PARSER:
        print("  ✅ 模板解析模块已加载")
    else:
        print("  ⚠️  模板解析模块未加载（需要 template_parser.py）")
    print()
    
    # ===== 步骤1：选择输入方式 =====
    print("=" * 70)
    print("📄 步骤1：输入PPT大纲")
    print("=" * 70)
    
    menu_items = [
        "[1] 弹窗粘贴大纲（推荐 - 最方便）",
        "[2] 终端粘贴大纲",
        "[3] 使用内置示例（军事主题）",
        "[4] 导入JSON文件",
    ]
    
    if HAS_TEMPLATE_PARSER:
        menu_items.append("[5] 🆕 使用企业模板生成")
        menu_items.append("[6] 🆕 分析模板样式（仅查看）")
    
    choice = input(
        "\n请选择输入方式:\n" + "\n".join(menu_items) + "\n默认: 1 > "
    ).strip() or "1"
    
    json_data = None
    json_path = None
    
    if choice == "1":
        # GUI弹窗输入模式
        print("\n📝 正在打开输入窗口...")
        text_content = get_text_from_gui()
        
        if not text_content:
            print("❌ 已取消或未输入内容")
            return
        
        # 尝试解析为JSON
        try:
            json_data = json.loads(text_content)
            print(f"\n✅ JSON格式解析成功")
        except json.JSONDecodeError:
            # 不是JSON，尝试解析为大纲文本
            print(f"\n📝 检测到纯文本大纲，正在转换...")
            json_data = parse_outline_to_json(text_content)
            if json_data:
                # 统计content_image数量
                image_count = sum(1 for s in json_data.get('slides', []) if s.get('type') == 'content_image')
                print(f"✅ 大纲转换成功，共 {len(json_data.get('slides', []))} 页幻灯片")
                print(f"   📸 其中 {image_count} 页支持AI图片生成")
            else:
                print("❌ 大纲解析失败")
                return
    
    elif choice == "2":
        # 终端粘贴文本模式
        print("\n" + "-" * 50)
        print("📝 请粘贴大纲文本（支持多行），输入完成后按两次回车结束：")
        print("-" * 50)
        
        lines = []
        empty_count = 0
        while True:
            try:
                line = input()
                if line == "":
                    empty_count += 1
                    if empty_count >= 2:
                        break
                    lines.append(line)
                else:
                    empty_count = 0
                    lines.append(line)
            except EOFError:
                break
        
        text_content = "\n".join(lines).strip()
        
        if not text_content:
            print("❌ 未输入任何内容")
            return
        
        # 尝试解析为JSON
        try:
            json_data = json.loads(text_content)
            print(f"\n✅ JSON格式解析成功")
        except json.JSONDecodeError:
            # 不是JSON，尝试解析为大纲文本
            print(f"\n📝 检测到纯文本大纲，正在转换...")
            json_data = parse_outline_to_json(text_content)
            if json_data:
                image_count = sum(1 for s in json_data.get('slides', []) if s.get('type') == 'content_image')
                print(f"✅ 大纲转换成功，共 {len(json_data.get('slides', []))} 页幻灯片")
                print(f"   📸 其中 {image_count} 页支持AI图片生成")
            else:
                print("❌ 大纲解析失败")
                return
    
    elif choice == "3":
        # 内置示例
        json_path = "example_simple.json"
        json_data = {
            "metadata": {"title": "示例演示", "theme": "military_solemn", "version": "3.9", "total_slides": 3},
            "slides": [
                {"type": "cover", "title": "PPT生成器v3.9测试", "subtitle": "完美版演示", "slogan": "GUI弹窗输入"},
                {"type": "section", "title": "核心改进"},
                {"type": "content_image", "title": "新功能展示", "bullets": ["GUI弹窗输入", "自动支持AI图片"], "layout": "left_text_right_image", "image_desc": "新功能展示"},
                {"type": "ending", "title": "测试完成", "bullets": ["流程优化", "多主题支持", "智能避让"], "quote": "完美！"}
            ]
        }
        with open(json_path, 'w', encoding='utf-8') as f:
            json.dump(json_data, f, ensure_ascii=False, indent=2)
        
        print(f"\n✅ 使用内置示例: {json_path}")
    
    elif choice == "4":
        # JSON文件导入
        json_path = input("\n请输入JSON文件路径: ").strip()
        if not os.path.exists(json_path):
            print(f"❌ 文件不存在: {json_path}")
            return
        
        try:
            with open(json_path, 'r', encoding='utf-8') as f:
                json_data = json.load(f)
            print(f"\n✅ JSON文件读取成功")
        except Exception as e:
            print(f"❌ JSON解析失败: {e}")
            return
    
    elif choice == "5" and HAS_TEMPLATE_PARSER:
        # ===== 🆕 使用企业模板生成 =====
        print("\n" + "=" * 70)
        print("📋 使用企业模板生成PPT")
        print("=" * 70)
        
        # 选择模板文件
        print("\n请选择模板文件...")
        template_path = select_template_file()
        
        if not template_path or not os.path.exists(template_path):
            print("❌ 未选择模板文件或文件不存在")
            return
        
        print(f"✅ 模板: {template_path}")
        
        # 分析模板
        print("\n📊 正在分析模板样式...")
        try:
            template_style = analyze_template(template_path)
        except Exception as e:
            print(f"❌ 模板分析失败: {e}")
            return
        
        # 输入大纲内容
        print("\n📝 请输入大纲内容...")
        text_content = get_text_from_gui()
        
        if not text_content:
            print("❌ 已取消或未输入内容")
            return
        
        # 解析大纲
        try:
            json_data = json.loads(text_content)
            print(f"✅ JSON格式解析成功")
        except json.JSONDecodeError:
            print(f"📝 检测到纯文本大纲，正在转换...")
            json_data = parse_outline_to_json(text_content)
            if not json_data:
                print("❌ 大纲解析失败")
                return
        
        # 选择生成模式
        print("\n生成模式：")
        print("  [1] 克隆模式 - 提取模板样式生成新PPT（推荐）")
        print("  [2] 填充模式 - 尝试复制模板页面并填充")
        mode_choice = input("选择 (默认: 1): ").strip() or "1"
        mode = 'fill' if mode_choice == '2' else 'clone'
        
        # 输出文件
        output_path = input("\n输出文件名 (默认: template_output.pptx): ").strip() or "template_output.pptx"
        if not output_path.endswith('.pptx'):
            output_path += '.pptx'
        
        # 使用模板生成
        try:
            generate_from_template(template_path, json_data, output_path, mode=mode)
            print("=" * 70)
            print(f"✅ 完成！文件已保存到: {output_path}")
            print("=" * 70)
        except Exception as e:
            print(f"❌ 生成失败: {e}")
        
        return  # 模板模式结束后直接返回
    
    elif choice == "6" and HAS_TEMPLATE_PARSER:
        # ===== 🆕 分析模板样式（仅查看） =====
        print("\n" + "=" * 70)
        print("🔍 模板样式分析")
        print("=" * 70)
        
        print("\n请选择模板文件...")
        template_path = select_template_file()
        
        if not template_path or not os.path.exists(template_path):
            print("❌ 未选择模板文件或文件不存在")
            return
        
        try:
            analyze_template(template_path)
            
            # 询问是否导出主题配置
            export_choice = input("\n是否导出为自定义主题配置？[y/N]: ").strip().lower()
            if export_choice == 'y':
                theme = get_theme_from_template(template_path)
                print("\n📋 可复制以下主题配置到代码中：")
                print("-" * 50)
                print("'custom_template': {")
                print(f"    'name': '自定义模板主题',")
                for key, value in theme.items():
                    if key == 'name':
                        continue
                    if hasattr(value, '__class__') and value.__class__.__name__ == 'RGBColor':
                        # 打印RGBColor信息
                        print(f"    '{key}': RGBColor({value.red}, {value.green}, {value.blue}),")
                    else:
                        print(f"    '{key}': {repr(value)},")
                print("}")
                print("-" * 50)
        except Exception as e:
            print(f"❌ 分析失败: {e}")
        
        return  # 仅分析模式结束后直接返回
    
    else:
        print("❌ 无效选择")
        return
    
    # ===== 步骤2：提取图片任务 =====
    image_tasks = extract_image_prompts_from_json(json_data)
    
    if image_tasks:
        print(f"\n📊 检测到 {len(image_tasks)} 个图片需要生成")
        for task in image_tasks:
            print(f"  - [{task.get('title', task['desc'])}]: {os.path.basename(task['path'])}")
    else:
        print("\n⚠️  JSON中没有图片配置（没有type为content_image的幻灯片）")
    
    # ===== 步骤3：询问是否下载图片 =====
    if image_tasks:
        print("\n" + "=" * 70)
        print("🖼️  步骤2：AI图片生成")
        print("=" * 70)
        
        download_choice = input(
            "\n是否生成/下载图片？\n"
            "[1] 是 (默认AI生成 - 推荐)\n"
            "[2] AI+Unsplash混合\n"
            "[3] 仅Unsplash\n"
            "[4] 否 (使用占位图)\n"
            "默认: 1 > "
        ).strip() or "1"
        
        unsplash_key = None
        siliconflow_key = None
        
        # 默认使用硅基流动AI
        if download_choice in ["1", "2"]:
            default_key = "sk-hkorbdmtxbynvelrlogfycsnckwlqqvjpdykodtrlxorxwcz"
            use_default = input(f"\n使用默认硅基流动Key? [Y/n]: ").strip().lower()
            
            if use_default != 'n':
                siliconflow_key = default_key
                print("✅ 使用默认AI Key")
            else:
                siliconflow_key = input("请输入硅基流动API Key: ").strip()
                if not siliconflow_key:
                    siliconflow_key = default_key
                    print("✅ 使用默认AI Key")
        
        if download_choice in ["2", "3"]:
            unsplash_key = input("\n请输入Unsplash Access Key (可选，回车跳过): ").strip()
            if not unsplash_key:
                print("⚠️  未提供Unsplash Key，仅使用AI生成")
                unsplash_key = None
        
        # 执行下载
        if download_choice in ["1", "2", "3"]:
            success_paths = download_images_from_json(image_tasks, unsplash_key, siliconflow_key)
            
            # 更新JSON中的图片路径（使用绝对路径）
            print("\n🔄 同步图片路径到JSON...")
            slide_index = 0
            for slide in json_data.get('slides', []):
                if slide.get('type') == 'content_image':
                    if slide_index < len(image_tasks):
                        # 使用绝对路径
                        abs_path = image_tasks[slide_index]['path']
                        slide['image'] = abs_path
                        
                        # 检查文件是否存在
                        if os.path.exists(abs_path):
                            print(f"  ✅ 更新: {os.path.basename(abs_path)} (存在)")
                        else:
                            print(f"  ⚠️ 更新: {os.path.basename(abs_path)} (不存在)")
                        slide_index += 1
            print("✅ 路径同步完成\n")
        else:
            print("\n✅ 将使用占位图\n")
    
    # ===== 步骤4：选择主题 =====
    print("=" * 70)
    print("🎨 步骤3：选择PPT主题")
    print("=" * 70)
    
    print("\n可用主题：")
    for i, (key, theme) in enumerate(AutoPPTGeneratorV3.THEMES.items(), 1):
        print(f"  [{i}] {theme['name']} ({key})")
    
    theme_choice = input("\n选择主题编号 (默认: 1-军事庄重): ").strip() or "1"
    theme_list = list(AutoPPTGeneratorV3.THEMES.keys())
    
    try:
        theme_index = int(theme_choice) - 1
        if 0 <= theme_index < len(theme_list):
            theme = theme_list[theme_index]
        else:
            theme = 'military_solemn'
    except:
        # 尝试作为主题名称
        if theme_choice in AutoPPTGeneratorV3.THEMES:
            theme = theme_choice
        else:
            # 从JSON中读取
            theme = json_data.get('metadata', {}).get('theme', 'military_solemn')
    
    # ===== 步骤5：生成PPT =====
    print("\n" + "=" * 70)
    print("📝 步骤4：生成PPT")
    print("=" * 70)
    
    output_path = input("\n输出文件名 (默认: output.pptx): ").strip() or "output.pptx"
    if not output_path.endswith('.pptx'):
        output_path += '.pptx'
    
    print()
    generator = AutoPPTGeneratorV3(theme=theme)
    # 直接传入更新后的json_data（包含正确的图片路径）
    generator.generate_from_json(json_data, output_path)
    
    print("=" * 70)
    print(f"✅ 完成！文件已保存到: {output_path}")
    print("=" * 70)


if __name__ == '__main__':
    main()
