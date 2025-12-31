#!/usr/bin/env python3
"""
PPT模板解析器模块 v1.0
功能：
1. 从现有PPT模板中提取样式要素（颜色、字体、布局）
2. 识别模板中的占位符类型
3. 支持基于模板生成新PPT（两种模式）

作者：AI资源指挥官
版本：1.0
更新：2025-12-31
"""

import os
import re
from copy import deepcopy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE


# ========================================================================
# 模板样式提取器
# ========================================================================

class TemplateStyleExtractor:
    """
    从PPT模板中提取样式要素
    提取内容：主题色、强调色、字体、布局等
    """
    
    def __init__(self, template_path):
        """
        初始化提取器
        
        Args:
            template_path: 模板PPT文件路径
        """
        if not os.path.exists(template_path):
            raise FileNotFoundError(f"模板文件不存在: {template_path}")
        
        self.template_path = template_path
        self.prs = Presentation(template_path)
        self.extracted_style = None
        
    def extract_all(self):
        """
        提取模板的所有样式信息
        
        Returns:
            dict: 包含颜色、字体、布局等的完整样式配置
        """
        self.extracted_style = {
            'slide_size': self._extract_slide_size(),
            'colors': self._extract_colors(),
            'fonts': self._extract_fonts(),
            'layouts': self._extract_layouts(),
            'backgrounds': self._extract_backgrounds(),
            'slide_masters': self._extract_slide_masters_info(),
        }
        
        return self.extracted_style
    
    def _extract_slide_size(self):
        """提取幻灯片尺寸"""
        return {
            'width': self.prs.slide_width,
            'height': self.prs.slide_height,
            'width_inches': self.prs.slide_width.inches,
            'height_inches': self.prs.slide_height.inches,
        }
    
    def _extract_colors(self):
        """
        从模板中提取主要使用的颜色
        分析所有形状和文本的颜色使用情况
        """
        colors = {
            'fill_colors': [],      # 填充色
            'text_colors': [],      # 文字颜色
            'line_colors': [],      # 边框颜色
            'primary': None,        # 推断的主色
            'accent': None,         # 推断的强调色
            'text': None,           # 推断的文字色
            'background': None,     # 推断的背景色
        }
        
        fill_color_count = {}
        text_color_count = {}
        bg_colors = []
        
        for slide in self.prs.slides:
            # 分析背景
            try:
                bg = slide.background
                if bg.fill and bg.fill.type is not None:
                    try:
                        if bg.fill.fore_color and bg.fill.fore_color.type is not None:
                            rgb = bg.fill.fore_color.rgb
                            if rgb:
                                bg_tuple = self._rgb_to_tuple(rgb)
                                if bg_tuple:
                                    bg_colors.append(bg_tuple)
                    except:
                        pass
            except:
                pass
            
            # 分析形状
            for shape in slide.shapes:
                # 填充色 - 更健壮的检测方式
                try:
                    if hasattr(shape, 'fill') and shape.fill:
                        fill = shape.fill
                        if fill.type is not None and hasattr(fill, 'fore_color'):
                            try:
                                if fill.fore_color and fill.fore_color.type is not None:
                                    rgb = fill.fore_color.rgb
                                    if rgb:
                                        rgb_tuple = self._rgb_to_tuple(rgb)
                                        if rgb_tuple:
                                            colors['fill_colors'].append(rgb_tuple)
                                            fill_color_count[rgb_tuple] = fill_color_count.get(rgb_tuple, 0) + 1
                            except:
                                pass
                except:
                    pass
                
                # 文字颜色 - 更健壮的检测方式
                try:
                    if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            # 检查段落级别字体颜色
                            try:
                                if paragraph.font and paragraph.font.color:
                                    if paragraph.font.color.type is not None:
                                        rgb = paragraph.font.color.rgb
                                        if rgb:
                                            rgb_tuple = self._rgb_to_tuple(rgb)
                                            if rgb_tuple:
                                                colors['text_colors'].append(rgb_tuple)
                                                text_color_count[rgb_tuple] = text_color_count.get(rgb_tuple, 0) + 1
                            except:
                                pass
                            
                            # 检查run级别字体颜色
                            for run in paragraph.runs:
                                try:
                                    if run.font and run.font.color:
                                        if run.font.color.type is not None:
                                            rgb = run.font.color.rgb
                                            if rgb:
                                                rgb_tuple = self._rgb_to_tuple(rgb)
                                                if rgb_tuple:
                                                    colors['text_colors'].append(rgb_tuple)
                                                    text_color_count[rgb_tuple] = text_color_count.get(rgb_tuple, 0) + 1
                                except:
                                    pass
                except:
                    pass
        
        # 保存背景色
        if bg_colors:
            colors['background'] = bg_colors[0]
        
        # 去重
        colors['fill_colors'] = list(set(colors['fill_colors']))
        colors['text_colors'] = list(set(colors['text_colors']))
        
        # 推断主要颜色
        if fill_color_count:
            # 排除白色和接近白色的颜色作为主色
            valid_fills = {k: v for k, v in fill_color_count.items() 
                          if sum(k) < 700}  # 排除接近白色的
            if valid_fills:
                colors['primary'] = max(valid_fills, key=valid_fills.get)
        
        if text_color_count:
            # 找最常用的深色文字
            dark_texts = {k: v for k, v in text_color_count.items() 
                         if sum(k) < 400}  # 深色文字
            if dark_texts:
                colors['text'] = max(dark_texts, key=dark_texts.get)
            
            # 找强调色（非黑非白的鲜艳颜色）
            accent_candidates = {k: v for k, v in text_color_count.items() 
                                if 150 < sum(k) < 600 and 
                                max(k) - min(k) > 50}  # 有色彩的
            if accent_candidates:
                colors['accent'] = max(accent_candidates, key=accent_candidates.get)
        
        return colors
    
    def _extract_fonts(self):
        """提取模板中使用的字体信息"""
        fonts = {
            'title_fonts': [],
            'body_fonts': [],
            'all_fonts': set(),
            'title_size': None,
            'body_size': None,
        }
        
        title_sizes = []
        body_sizes = []
        
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        try:
                            if run.font.name:
                                fonts['all_fonts'].add(run.font.name)
                            
                            # 根据字号判断是标题还是正文
                            if run.font.size:
                                size_pt = run.font.size.pt
                                if size_pt >= 24:
                                    title_sizes.append(size_pt)
                                    if run.font.name:
                                        fonts['title_fonts'].append(run.font.name)
                                else:
                                    body_sizes.append(size_pt)
                                    if run.font.name:
                                        fonts['body_fonts'].append(run.font.name)
                        except:
                            pass
        
        # 转换为列表
        fonts['all_fonts'] = list(fonts['all_fonts'])
        
        # 计算平均字号
        if title_sizes:
            fonts['title_size'] = sum(title_sizes) / len(title_sizes)
        if body_sizes:
            fonts['body_size'] = sum(body_sizes) / len(body_sizes)
        
        return fonts
    
    def _extract_layouts(self):
        """提取模板的布局信息"""
        layouts = []
        
        for idx, slide in enumerate(self.prs.slides):
            slide_layout = {
                'index': idx,
                'shapes': [],
                'has_title': False,
                'has_content': False,
                'has_image_placeholder': False,
            }
            
            for shape in slide.shapes:
                shape_info = {
                    'type': str(shape.shape_type),
                    'left': shape.left,
                    'top': shape.top,
                    'width': shape.width,
                    'height': shape.height,
                    'name': shape.name,
                    'has_text': shape.has_text_frame,
                }
                
                # 判断是否为标题
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if '标题' in shape.name.lower() or 'title' in shape.name.lower():
                        slide_layout['has_title'] = True
                    elif text:
                        slide_layout['has_content'] = True
                
                # 判断是否有图片占位符
                if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                    slide_layout['has_image_placeholder'] = True
                
                slide_layout['shapes'].append(shape_info)
            
            layouts.append(slide_layout)
        
        return layouts
    
    def _extract_backgrounds(self):
        """提取各页面的背景设置"""
        backgrounds = []
        
        for idx, slide in enumerate(self.prs.slides):
            bg_info = {'index': idx, 'type': 'unknown'}
            
            try:
                bg = slide.background
                fill = bg.fill
                
                if fill.type == 1:  # 纯色
                    bg_info['type'] = 'solid'
                    if fill.fore_color.rgb:
                        bg_info['color'] = self._rgb_to_tuple(fill.fore_color.rgb)
                elif fill.type == 2:  # 渐变
                    bg_info['type'] = 'gradient'
                elif fill.type == 3:  # 图片
                    bg_info['type'] = 'picture'
                elif fill.type == 4:  # 图案
                    bg_info['type'] = 'pattern'
            except:
                pass
            
            backgrounds.append(bg_info)
        
        return backgrounds
    
    def _extract_slide_masters_info(self):
        """提取母版信息"""
        masters = []
        
        for master in self.prs.slide_masters:
            master_info = {
                'layouts_count': len(master.slide_layouts),
                'layout_names': [layout.name for layout in master.slide_layouts]
            }
            masters.append(master_info)
        
        return masters
    
    def _rgb_to_tuple(self, rgb):
        """将RGBColor转换为元组"""
        if isinstance(rgb, RGBColor):
            return (rgb.red, rgb.green, rgb.blue)
        return None
    
    def get_theme_config(self):
        """
        将提取的样式转换为生成器可用的主题配置
        
        Returns:
            dict: 可直接用于AutoPPTGeneratorV3的主题配置
        """
        if not self.extracted_style:
            self.extract_all()
        
        colors = self.extracted_style['colors']
        
        # 构建主题配置
        theme = {
            'name': '自定义模板主题',
            'primary': self._tuple_to_rgb(colors.get('primary')) or RGBColor(26, 35, 126),
            'accent': self._tuple_to_rgb(colors.get('accent')) or RGBColor(213, 0, 0),
            'text': self._tuple_to_rgb(colors.get('text')) or RGBColor(33, 33, 33),
            'bg': RGBColor(250, 250, 250),
            'quote': self._tuple_to_rgb(colors.get('accent')) or RGBColor(0, 150, 136),
            'chart': self._tuple_to_rgb(colors.get('primary')) or RGBColor(63, 81, 181),
        }
        
        return theme
    
    def _tuple_to_rgb(self, color_tuple):
        """将颜色元组转换为RGBColor"""
        if color_tuple and len(color_tuple) == 3:
            return RGBColor(color_tuple[0], color_tuple[1], color_tuple[2])
        return None
    
    def print_summary(self):
        """打印提取的样式摘要"""
        if not self.extracted_style:
            self.extract_all()
        
        print("\n" + "="*60)
        print("📊 模板样式分析报告")
        print("="*60)
        
        # 尺寸
        size = self.extracted_style['slide_size']
        print(f"\n📐 幻灯片尺寸: {size['width_inches']:.2f}\" x {size['height_inches']:.2f}\"")
        
        # 颜色
        colors = self.extracted_style['colors']
        print(f"\n🎨 颜色分析:")
        if colors.get('primary'):
            print(f"   主色: RGB{colors['primary']}")
        if colors.get('accent'):
            print(f"   强调色: RGB{colors['accent']}")
        if colors.get('text'):
            print(f"   文字色: RGB{colors['text']}")
        print(f"   填充色种类: {len(colors['fill_colors'])}种")
        print(f"   文字色种类: {len(colors['text_colors'])}种")
        
        # 字体
        fonts = self.extracted_style['fonts']
        print(f"\n🔤 字体分析:")
        print(f"   使用的字体: {', '.join(fonts['all_fonts'][:5]) if fonts['all_fonts'] else '未检测到'}")
        if fonts['title_size']:
            print(f"   标题字号: 约{fonts['title_size']:.1f}pt")
        if fonts['body_size']:
            print(f"   正文字号: 约{fonts['body_size']:.1f}pt")
        
        # 布局
        layouts = self.extracted_style['layouts']
        print(f"\n📄 页面分析:")
        print(f"   总页数: {len(layouts)}")
        for layout in layouts[:5]:
            print(f"   第{layout['index']+1}页: {len(layout['shapes'])}个形状")
        
        # 母版
        masters = self.extracted_style['slide_masters']
        print(f"\n🎭 母版信息:")
        for i, master in enumerate(masters):
            print(f"   母版{i+1}: {master['layouts_count']}种布局")
            if master['layout_names']:
                print(f"   布局: {', '.join(master['layout_names'][:4])}...")
        
        print("\n" + "="*60)


# ========================================================================
# 基于模板的PPT生成器
# ========================================================================

class TemplateBasedGenerator:
    """
    基于模板生成PPT的生成器
    支持两种模式：
    1. 占位符填充模式：直接在模板的占位符中填充内容
    2. 样式克隆模式：使用模板样式，但重新生成页面结构
    """
    
    def __init__(self, template_path):
        """
        初始化生成器
        
        Args:
            template_path: 模板PPT文件路径
        """
        if not os.path.exists(template_path):
            raise FileNotFoundError(f"模板文件不存在: {template_path}")
        
        self.template_path = template_path
        self.template_prs = Presentation(template_path)
        
        # 提取模板样式
        self.extractor = TemplateStyleExtractor(template_path)
        self.style = self.extractor.extract_all()
        self.theme = self.extractor.get_theme_config()
        
        # 用于生成的演示文稿
        self.prs = None
        self.slide_index = 0
        
        print(f"✅ 模板加载成功: {template_path}")
        print(f"📄 模板包含 {len(self.template_prs.slides)} 个页面")
    
    def generate_from_json(self, json_path_or_data, output_path, mode='clone'):
        """
        基于模板从JSON生成PPT
        
        Args:
            json_path_or_data: JSON文件路径或字典数据
            output_path: 输出路径
            mode: 生成模式
                  'clone' - 克隆模板样式生成新结构
                  'fill' - 使用模板页面填充内容
        """
        # 加载JSON数据
        if isinstance(json_path_or_data, dict):
            data = json_path_or_data
        else:
            import json
            with open(json_path_or_data, 'r', encoding='utf-8') as f:
                data = json.load(f)
        
        slides_data = data.get('slides', [])
        
        print(f"\n{'='*60}")
        print(f"🚀 开始基于模板生成PPT (模式: {mode})")
        print(f"{'='*60}\n")
        
        if mode == 'fill':
            self._generate_fill_mode(slides_data, output_path)
        else:
            self._generate_clone_mode(slides_data, output_path)
        
        print(f"\n{'='*60}")
        print(f"✅ PPT生成成功！")
        print(f"📊 总页数: {len(self.prs.slides)}")
        print(f"📁 输出路径: {output_path}")
        print(f"{'='*60}\n")
    
    def _generate_fill_mode(self, slides_data, output_path):
        """
        填充模式：复制模板页面，填充内容
        
        适用于模板有明确占位符结构的情况
        """
        # 创建模板的副本
        self.prs = Presentation(self.template_path)
        
        # 获取模板中的页面类型映射
        template_slides = self._analyze_template_slides()
        
        # 清空现有内容（保留第一页作为模板）
        # 注意：python-pptx不支持直接删除幻灯片，我们需要另一种方式
        
        # 创建新的演示文稿，但使用模板的母版
        self.prs = Presentation()
        self.prs.slide_width = self.template_prs.slide_width
        self.prs.slide_height = self.template_prs.slide_height
        
        for slide_data in slides_data:
            slide_type = slide_data.get('type')
            
            # 选择最匹配的模板页面
            template_slide = self._find_matching_template_slide(slide_type, template_slides)
            
            if template_slide:
                # 复制模板页面并填充
                new_slide = self._copy_and_fill_slide(template_slide, slide_data)
            else:
                # 没有匹配模板，使用克隆模式生成
                self._create_slide_with_style(slide_data)
            
            self.slide_index += 1
        
        self.prs.save(output_path)
    
    def _generate_clone_mode(self, slides_data, output_path):
        """
        克隆模式：使用模板样式生成新结构
        
        这是更通用的方式，提取模板的颜色/字体等样式后生成
        """
        # 创建新的演示文稿
        self.prs = Presentation()
        self.prs.slide_width = self.template_prs.slide_width
        self.prs.slide_height = self.template_prs.slide_height
        
        for slide_data in slides_data:
            self._create_slide_with_style(slide_data)
            self.slide_index += 1
        
        self.prs.save(output_path)
    
    def _analyze_template_slides(self):
        """
        分析模板中各页面的类型
        
        Returns:
            list: 每页的类型信息
        """
        slides_info = []
        
        for idx, slide in enumerate(self.template_prs.slides):
            info = {
                'index': idx,
                'slide': slide,
                'type': 'content',  # 默认类型
                'has_title': False,
                'has_subtitle': False,
                'has_content': False,
                'has_image': False,
                'placeholder_types': [],
            }
            
            for shape in slide.shapes:
                # 检查形状类型和内容
                shape_name = shape.name.lower()
                
                if 'title' in shape_name or '标题' in shape_name:
                    info['has_title'] = True
                elif 'subtitle' in shape_name or '副标题' in shape_name:
                    info['has_subtitle'] = True
                
                if shape.has_text_frame:
                    text = shape.text_frame.text.lower()
                    if '封面' in text or 'cover' in text:
                        info['type'] = 'cover'
                    elif '目录' in text or 'contents' in text:
                        info['type'] = 'section'
                    elif '结束' in text or '谢谢' in text or 'thank' in text:
                        info['type'] = 'ending'
                    elif '图表' in text or 'chart' in text:
                        info['type'] = 'chart'
                    else:
                        info['has_content'] = True
                
                # 检查是否有图片
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    info['has_image'] = True
            
            # 根据特征推断类型
            if idx == 0 and info['has_title'] and info['has_subtitle']:
                info['type'] = 'cover'
            elif idx == len(self.template_prs.slides) - 1:
                info['type'] = 'ending'
            elif info['has_image'] and info['has_content']:
                info['type'] = 'content_image'
            
            slides_info.append(info)
            print(f"  📄 第{idx+1}页: 类型={info['type']}")
        
        return slides_info
    
    def _find_matching_template_slide(self, slide_type, template_slides):
        """
        根据类型找到匹配的模板页面
        """
        for slide_info in template_slides:
            if slide_info['type'] == slide_type:
                return slide_info['slide']
        
        # 没找到完全匹配的，返回通用内容页
        for slide_info in template_slides:
            if slide_info['type'] == 'content':
                return slide_info['slide']
        
        return None
    
    def _copy_and_fill_slide(self, template_slide, data):
        """
        复制模板页面并填充内容
        """
        # 添加空白页
        layout = self.prs.slide_layouts[6]  # 空白布局
        new_slide = self.prs.slides.add_slide(layout)
        
        # 复制背景
        try:
            bg = template_slide.background
            new_bg = new_slide.background
            if bg.fill.type == 1:  # 纯色
                new_bg.fill.solid()
                new_bg.fill.fore_color.rgb = bg.fill.fore_color.rgb
        except:
            pass
        
        # 复制并填充形状
        title_text = data.get('title', '')
        subtitle_text = data.get('subtitle', '')
        bullets = data.get('bullets', [])
        
        for shape in template_slide.shapes:
            try:
                # 复制形状的位置和尺寸
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                
                shape_name = shape.name.lower()
                
                if shape.has_text_frame:
                    # 创建文本框
                    new_shape = new_slide.shapes.add_textbox(left, top, width, height)
                    new_tf = new_shape.text_frame
                    
                    # 填充内容
                    if 'title' in shape_name or '标题' in shape_name:
                        new_tf.text = title_text
                    elif 'subtitle' in shape_name or '副标题' in shape_name:
                        new_tf.text = subtitle_text
                    elif bullets:
                        # 填充要点
                        for i, bullet in enumerate(bullets):
                            if i == 0:
                                new_tf.paragraphs[0].text = bullet
                            else:
                                p = new_tf.add_paragraph()
                                p.text = bullet
                    
                    # 复制样式
                    self._copy_text_style(shape.text_frame, new_tf)
                    
            except Exception as e:
                print(f"  ⚠️ 复制形状时出错: {e}")
        
        return new_slide
    
    def _copy_text_style(self, source_tf, target_tf):
        """
        复制文本框的样式
        """
        try:
            for i, src_para in enumerate(source_tf.paragraphs):
                if i >= len(target_tf.paragraphs):
                    break
                
                tgt_para = target_tf.paragraphs[i]
                tgt_para.alignment = src_para.alignment
                
                for j, src_run in enumerate(src_para.runs):
                    if j < len(tgt_para.runs):
                        tgt_run = tgt_para.runs[j]
                    else:
                        continue
                    
                    try:
                        if src_run.font.size:
                            tgt_run.font.size = src_run.font.size
                        if src_run.font.bold is not None:
                            tgt_run.font.bold = src_run.font.bold
                        if src_run.font.color.rgb:
                            tgt_run.font.color.rgb = src_run.font.color.rgb
                        if src_run.font.name:
                            tgt_run.font.name = src_run.font.name
                    except:
                        pass
        except:
            pass
    
    def _create_slide_with_style(self, data):
        """
        使用提取的样式创建新页面
        """
        slide_type = data.get('type')
        
        if slide_type == 'cover':
            self._create_cover_slide(data)
        elif slide_type == 'section':
            self._create_section_slide(data)
        elif slide_type == 'content_image':
            self._create_content_slide(data)
        elif slide_type == 'chart':
            self._create_chart_slide(data)
        elif slide_type == 'ending':
            self._create_ending_slide(data)
        else:
            self._create_content_slide(data)
    
    def _create_cover_slide(self, data):
        """创建封面页（使用模板样式）"""
        layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(layout)
        
        # 背景
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.theme['primary']
        
        # 主标题
        title_text = data.get('title', '')
        title_box = slide.shapes.add_textbox(
            Inches(0.3), Inches(1.5), Inches(9.4), Inches(1.5)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.text = title_text
        
        p = tf.paragraphs[0]
        p.font.size = Pt(40 if len(title_text) <= 15 else 32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # 副标题
        if data.get('subtitle'):
            subtitle_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(3.2), Inches(9), Inches(0.8)
            )
            tf = subtitle_box.text_frame
            tf.word_wrap = True
            tf.text = data['subtitle']
            
            p = tf.paragraphs[0]
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(230, 230, 230)
            p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def _create_section_slide(self, data):
        """创建章节页（使用模板样式）"""
        layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(layout)
        
        # 背景
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.theme.get('bg', RGBColor(250, 250, 250))
        
        # 装饰条
        deco = slide.shapes.add_shape(
            1, Inches(0), Inches(2.3),
            Inches(10), Inches(1)
        )
        deco.fill.solid()
        deco.fill.fore_color.rgb = self.theme['primary']
        deco.line.color.rgb = self.theme['primary']
        
        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.3), Inches(2.3), Inches(9.4), Inches(1)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.text = data.get('title', '')
        
        p = tf.paragraphs[0]
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def _create_content_slide(self, data):
        """创建内容页（使用模板样式）"""
        layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(layout)
        
        # 背景
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.theme.get('bg', RGBColor(250, 250, 250))
        
        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.3), Inches(0.3), Inches(9.4), Inches(0.8)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.text = data.get('title', '')
        
        p = tf.paragraphs[0]
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.theme['primary']
        
        # 内容
        if data.get('bullets'):
            content_box = slide.shapes.add_textbox(
                Inches(0.3), Inches(1.3), Inches(4.5), Inches(3.5)
            )
            tf = content_box.text_frame
            tf.word_wrap = True
            
            for i, bullet in enumerate(data['bullets']):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                p.text = f"• {bullet}"
                p.font.size = Pt(12)
                p.font.color.rgb = self.theme['text']
                p.space_before = Pt(4)
                p.space_after = Pt(4)
        
        # 图片占位符
        if data.get('image_path') or data.get('image_desc'):
            try:
                if data.get('image_path') and os.path.exists(data['image_path']):
                    slide.shapes.add_picture(
                        data['image_path'],
                        Inches(5.0), Inches(1.3),
                        width=Inches(4.5), height=Inches(3.5)
                    )
                else:
                    # 添加占位符
                    placeholder = slide.shapes.add_shape(
                        1, Inches(5.0), Inches(1.3),
                        Inches(4.5), Inches(3.5)
                    )
                    placeholder.fill.solid()
                    placeholder.fill.fore_color.rgb = RGBColor(240, 240, 240)
                    
                    # 添加描述文字
                    desc_box = slide.shapes.add_textbox(
                        Inches(5.2), Inches(2.5), Inches(4.1), Inches(1)
                    )
                    tf = desc_box.text_frame
                    tf.text = f"📷 {data.get('image_desc', '待添加图片')}"
                    tf.paragraphs[0].font.size = Pt(10)
                    tf.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
                    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            except Exception as e:
                print(f"  ⚠️ 添加图片时出错: {e}")
        
        return slide
    
    def _create_chart_slide(self, data):
        """创建图表页"""
        layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(layout)
        
        # 背景
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.theme.get('bg', RGBColor(250, 250, 250))
        
        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4), Inches(9), Inches(0.6)
        )
        tf = title_box.text_frame
        tf.text = data.get('title', '')
        
        p = tf.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.theme['primary']
        
        return slide
    
    def _create_ending_slide(self, data):
        """创建结束页"""
        layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(layout)
        
        # 背景
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.theme.get('bg', RGBColor(250, 250, 250))
        
        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.6), Inches(9), Inches(0.8)
        )
        tf = title_box.text_frame
        tf.text = data.get('title', '')
        
        p = tf.paragraphs[0]
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = self.theme['primary']
        
        return slide


# ========================================================================
# 便捷接口函数
# ========================================================================

def analyze_template(template_path):
    """
    分析PPT模板，打印样式报告
    
    Args:
        template_path: 模板文件路径
    
    Returns:
        dict: 提取的样式信息
    """
    extractor = TemplateStyleExtractor(template_path)
    style = extractor.extract_all()
    extractor.print_summary()
    return style


def get_theme_from_template(template_path):
    """
    从模板提取主题配置，可直接用于AutoPPTGeneratorV3
    
    Args:
        template_path: 模板文件路径
    
    Returns:
        dict: 主题配置字典
    """
    extractor = TemplateStyleExtractor(template_path)
    return extractor.get_theme_config()


def generate_from_template(template_path, json_data, output_path, mode='clone'):
    """
    基于模板生成PPT的快捷函数
    
    Args:
        template_path: 模板文件路径
        json_data: JSON数据（字典或文件路径）
        output_path: 输出文件路径
        mode: 'clone'(克隆样式) 或 'fill'(填充模板)
    
    Returns:
        str: 输出文件路径
    """
    generator = TemplateBasedGenerator(template_path)
    generator.generate_from_json(json_data, output_path, mode=mode)
    return output_path


# ========================================================================
# 测试/演示
# ========================================================================

if __name__ == '__main__':
    import sys
    
    print("\n" + "="*60)
    print("📋 PPT模板解析器 v1.0")
    print("="*60)
    
    if len(sys.argv) < 2:
        print("""
使用方法:
    python template_parser.py <模板路径>         # 分析模板
    python template_parser.py <模板> <json> <输出>  # 基于模板生成

示例:
    python template_parser.py company_template.pptx
    python template_parser.py template.pptx data.json output.pptx
        """)
        sys.exit(0)
    
    template_path = sys.argv[1]
    
    if len(sys.argv) == 2:
        # 仅分析模板
        analyze_template(template_path)
    elif len(sys.argv) >= 4:
        # 基于模板生成
        json_path = sys.argv[2]
        output_path = sys.argv[3]
        mode = sys.argv[4] if len(sys.argv) > 4 else 'clone'
        
        generate_from_template(template_path, json_path, output_path, mode)
